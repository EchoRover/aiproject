# Generate PowerPoint Presentation for AI Project

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

print("=" * 70)
print("GENERATING POWERPOINT PRESENTATION")
print("=" * 70)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
TITLE_COLOR = RGBColor(31, 78, 121)  # Dark blue
ACCENT_COLOR = RGBColor(52, 152, 219)  # Light blue

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add author and date
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    author_frame = author_box.text_frame
    author_frame.text = "Evan Tobias\nAI Final Year Project\nDecember 7, 2025"
    author_frame.paragraphs[0].font.size = Pt(18)
    author_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    # Add content
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for line in content_lines:
        p = text_frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(10)
    
    return slide

def add_image_slide(prs, title, image_path):
    """Add a slide with an image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.2), width=Inches(8))
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(5.5))
    left_frame = left_box.text_frame
    for line in left_content:
        p = left_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.2), Inches(4.25), Inches(5.5))
    right_frame = right_box.text_frame
    for line in right_content:
        p = right_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    
    return slide

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
print("\n📊 Creating Slide 1: Title...")
add_title_slide(prs, 
    "Machine Learning Algorithms Showcase",
    "Comprehensive Analysis on Energy Datasets")

# ============================================================================
# SLIDE 2: Agenda
# ============================================================================
print("📊 Creating Slide 2: Agenda...")
add_content_slide(prs, "Presentation Agenda", [
    "📋 Introduction & Project Overview",
    "📊 Datasets Description",
    "🔧 Methodology & Preprocessing",
    "🤖 Regression Algorithms (5 models)",
    "🎯 Classification with Logistic Regression",
    "🔍 Clustering with K-means",
    "📈 Results & Performance Comparison",
    "💡 Key Findings & Conclusions"
])

# ============================================================================
# SLIDE 3: Project Overview
# ============================================================================
print("📊 Creating Slide 3: Project Overview...")
add_content_slide(prs, "Project Overview", [
    "🎯 Objective: Demonstrate 7 ML algorithms on real-world data",
    "📚 Algorithms Implemented:",
    "   • Linear & Polynomial Regression",
    "   • Decision Trees & Random Forest",
    "   • Neural Networks (PyTorch)",
    "   • Logistic Regression (Classification)",
    "   • K-means Clustering",
    "🔬 Two Energy Datasets: ENB2012 & Energy Consumption",
    "📊 768 & 19,735 samples respectively"
])

# ============================================================================
# SLIDE 4: Datasets
# ============================================================================
print("📊 Creating Slide 4: Datasets...")
add_two_column_slide(prs, "Datasets Overview",
    [
        "ENB2012 Dataset",
        "• 768 samples",
        "• 8 features (building design)",
        "• Target: Heating/Cooling Load",
        "• Use: Regression analysis",
        "",
        "Features:",
        "• Relative Compactness",
        "• Surface & Wall Area",
        "• Overall Height",
        "• Glazing Area"
    ],
    [
        "Energy Consumption Dataset",
        "• 19,735 samples",
        "• 28 features (sensors)",
        "• Target: Appliances energy",
        "• Use: Classification & Clustering",
        "",
        "Features:",
        "• Temperature (multiple rooms)",
        "• Humidity (multiple rooms)",
        "• Pressure, Wind Speed",
        "• Energy consumption (Wh)"
    ])

# ============================================================================
# SLIDE 5: Correlation Analysis
# ============================================================================
print("📊 Creating Slide 5: Correlation Analysis...")
add_image_slide(prs, "Correlation Analysis", "images/01_correlation_heatmaps.png")

# ============================================================================
# SLIDE 6: Methodology
# ============================================================================
print("📊 Creating Slide 6: Methodology...")
add_content_slide(prs, "Methodology & Preprocessing", [
    "🧹 Data Cleaning: No missing values, no duplicates",
    "📏 Feature Scaling: Z-score normalization (μ=0, σ=1)",
    "🔍 Multicollinearity Check: VIF analysis",
    "✂️ Train-Test Split: 80/20 ratio, stratified for classification",
    "🎯 Feature Engineering: Binary target for classification",
    "🔧 Tools: scikit-learn, PyTorch, pandas, numpy"
])

# ============================================================================
# SLIDE 7: Linear Regression
# ============================================================================
print("📊 Creating Slide 7: Linear Regression...")
add_content_slide(prs, "Linear Regression", [
    "📐 Formula: ŷ = β₀ + β₁x₁ + β₂x₂ + ... + βₙxₙ",
    "🎯 Method: Ordinary Least Squares (OLS)",
    "📊 Loss Function: Mean Squared Error (MSE)",
    "",
    "✅ Results:",
    "   • R² Score: 0.9122 (91.22% variance explained)",
    "   • RMSE: 3.0254",
    "   • MAE: 2.1821",
    "",
    "💡 Strong baseline performance with linear relationships"
])

# ============================================================================
# SLIDE 8: Polynomial Regression
# ============================================================================
print("📊 Creating Slide 8: Polynomial Regression...")
add_content_slide(prs, "Polynomial Regression", [
    "📐 Extends linear regression with polynomial terms (degree=2)",
    "🔧 Includes interaction terms: x₁x₂, x₁², x₂²",
    "🎯 Captures non-linear relationships",
    "",
    "✅ Results:",
    "   • R² Score: 0.9938 (99.38% variance explained) 🌟",
    "   • RMSE: 0.8030",
    "   • MAE: 0.6042",
    "",
    "💡 Excellent performance, captures complex patterns"
])

# ============================================================================
# SLIDE 9: Decision Tree & Random Forest
# ============================================================================
print("📊 Creating Slide 9: Decision Tree & Random Forest...")
add_two_column_slide(prs, "Decision Tree & Random Forest",
    [
        "🌳 Decision Tree",
        "• Recursive partitioning",
        "• Non-parametric",
        "• Interpretable splits",
        "",
        "Results:",
        "• R²: 0.9883",
        "• RMSE: 1.1059",
        "• MAE: 0.7561",
        "",
        "Top Feature:",
        "Overall Height (58%)"
    ],
    [
        "🌲 Random Forest",
        "• Ensemble of 100 trees",
        "• Bootstrap + Feature bagging",
        "• Robust to overfitting",
        "",
        "Results:",
        "• R²: 0.9976 🏆 BEST!",
        "• RMSE: 0.4978",
        "• MAE: 0.3584",
        "",
        "Near-perfect predictions"
    ])

# ============================================================================
# SLIDE 10: Feature Importance
# ============================================================================
print("📊 Creating Slide 10: Feature Importance...")
add_image_slide(prs, "Feature Importance Analysis", "images/05_feature_importance.png")

# ============================================================================
# SLIDE 11: Neural Network
# ============================================================================
print("📊 Creating Slide 11: Neural Network...")
add_content_slide(prs, "Neural Network (PyTorch)", [
    "🧠 Architecture: 8 → 64 → 32 → 16 → 1",
    "⚡ Activation: ReLU for hidden layers",
    "🎯 Total Parameters: 3,201",
    "🔧 Optimizer: Adam (lr=0.001)",
    "📈 Training: 200 epochs, MSE Loss",
    "",
    "✅ Results:",
    "   • R² Score: 0.9683 (96.83%)",
    "   • RMSE: 1.8186",
    "   • MAE: 1.3031",
    "",
    "💡 Strong performance, demonstrates backpropagation"
])

# ============================================================================
# SLIDE 12: Neural Network Architecture
# ============================================================================
print("📊 Creating Slide 12: NN Architecture...")
add_image_slide(prs, "Neural Network Architecture", "images/04_neural_network_architecture.png")

# ============================================================================
# SLIDE 13: Regression Comparison
# ============================================================================
print("📊 Creating Slide 13: Regression Comparison...")
add_image_slide(prs, "Regression Models Comparison", "images/02_regression_comparison.png")

# ============================================================================
# SLIDE 14: All Metrics Comparison
# ============================================================================
print("📊 Creating Slide 14: All Metrics...")
add_image_slide(prs, "Comprehensive Metrics Comparison", "images/03_metrics_comparison.png")

# ============================================================================
# SLIDE 15: Logistic Regression (Classification)
# ============================================================================
print("📊 Creating Slide 15: Classification...")
add_content_slide(prs, "Logistic Regression - Classification", [
    "🎯 Task: Classify energy consumption (Low vs High)",
    "📐 Method: Sigmoid function for probability",
    "🔧 Binary threshold at median consumption",
    "",
    "✅ Results:",
    "   • Accuracy: 75.65%",
    "   • Precision: 73.69%",
    "   • Recall: 80.21%",
    "   • F1-Score: 76.81%",
    "   • AUC-ROC: 0.8329 🌟",
    "",
    "💡 Good discriminative ability, high recall"
])

# ============================================================================
# SLIDE 16: Confusion Matrix & ROC
# ============================================================================
print("📊 Creating Slide 16: Confusion Matrix...")
add_image_slide(prs, "Classification Performance", "images/06_confusion_matrix.png")

# ============================================================================
# SLIDE 17: ROC Curve
# ============================================================================
print("📊 Creating Slide 17: ROC Curve...")
add_image_slide(prs, "ROC Curve Analysis", "images/07_roc_curve.png")

# ============================================================================
# SLIDE 18: K-means Clustering
# ============================================================================
print("📊 Creating Slide 18: K-means Clustering...")
add_content_slide(prs, "K-means Clustering", [
    "🔍 Task: Discover consumption patterns",
    "🎯 Method: Partition data into k clusters",
    "📊 Optimal k Selection: Elbow + Silhouette methods",
    "",
    "✅ Results:",
    "   • Optimal k: 2 clusters",
    "   • Silhouette Score: 0.22",
    "   • Cluster 0: 9,894 samples (Low usage: 41 Wh)",
    "   • Cluster 1: 9,841 samples (High usage: 105 Wh)",
    "",
    "💡 Clear separation between usage patterns"
])

# ============================================================================
# SLIDE 19: Clustering Analysis
# ============================================================================
print("📊 Creating Slide 19: Clustering Analysis...")
add_image_slide(prs, "Clustering Analysis: Elbow & Silhouette", "images/08_clustering_analysis.png")

# ============================================================================
# SLIDE 20: Overall Summary
# ============================================================================
print("📊 Creating Slide 20: Overall Summary...")
add_image_slide(prs, "All Algorithms Performance Summary", "images/10_overall_summary.png")

# ============================================================================
# SLIDE 21: Key Findings
# ============================================================================
print("📊 Creating Slide 21: Key Findings...")
add_content_slide(prs, "Key Findings", [
    "🏆 Best Regression: Random Forest (R² = 0.9976)",
    "🧠 Neural Network: Excellent performance (R² = 0.9683)",
    "🎯 Classification: Strong results (75.65% accuracy, AUC = 0.83)",
    "🔍 Clustering: Identified 2 clear consumption patterns",
    "📊 All models achieved >91% R² for regression",
    "💡 Ensemble methods (Random Forest) excel",
    "⚡ Deep learning competitive with proper tuning",
    "🔧 Feature importance: Overall Height most critical (58%)"
])

# ============================================================================
# SLIDE 22: Conclusions
# ============================================================================
print("📊 Creating Slide 22: Conclusions...")
add_content_slide(prs, "Conclusions & Future Work", [
    "✅ Successfully implemented 7 ML algorithms",
    "✅ Comprehensive evaluation across multiple metrics",
    "✅ Real-world energy datasets with practical insights",
    "✅ Demonstrated supervised & unsupervised learning",
    "",
    "🔮 Future Work:",
    "   • LSTM/GRU for time-series forecasting",
    "   • Hyperparameter optimization (Grid/Bayesian search)",
    "   • Ensemble stacking for improved performance",
    "   • Real-time deployment as web service",
    "   • Explainable AI (SHAP/LIME) for interpretability"
])

# ============================================================================
# SLIDE 23: Thank You
# ============================================================================
print("📊 Creating Slide 23: Thank You...")
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

# Thank you text
thank_you_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
thank_you_frame = thank_you_box.text_frame
thank_you_frame.text = "Thank You!\n\nQuestions?"
thank_you_frame.paragraphs[0].font.size = Pt(54)
thank_you_frame.paragraphs[0].font.bold = True
thank_you_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
thank_you_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Contact info
contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
contact_frame = contact_box.text_frame
contact_frame.text = "Evan Tobias | AI Final Year Project | December 2025"
contact_frame.paragraphs[0].font.size = Pt(18)
contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================================
# Save Presentation
# ============================================================================
output_file = "AI_Project_Presentation.pptx"
prs.save(output_file)

print("\n" + "=" * 70)
print("✅ POWERPOINT PRESENTATION CREATED SUCCESSFULLY!")
print("=" * 70)
print(f"\n📁 File: {output_file}")
print(f"📊 Total Slides: {len(prs.slides)}")
print("\n📋 Slide Breakdown:")
print("  1. Title Slide")
print("  2. Agenda")
print("  3. Project Overview")
print("  4. Datasets Description")
print("  5. Correlation Analysis (Image)")
print("  6. Methodology & Preprocessing")
print("  7. Linear Regression")
print("  8. Polynomial Regression")
print("  9. Decision Tree & Random Forest")
print(" 10. Feature Importance (Image)")
print(" 11. Neural Network Details")
print(" 12. Neural Network Architecture (Image)")
print(" 13. Regression Comparison (Image)")
print(" 14. All Metrics Comparison (Image)")
print(" 15. Logistic Regression Classification")
print(" 16. Confusion Matrix (Image)")
print(" 17. ROC Curve (Image)")
print(" 18. K-means Clustering")
print(" 19. Clustering Analysis (Image)")
print(" 20. Overall Summary (Image)")
print(" 21. Key Findings")
print(" 22. Conclusions & Future Work")
print(" 23. Thank You / Q&A")
print("\n💡 Ready to present! Duration: ~15-20 minutes")
print("=" * 70)
